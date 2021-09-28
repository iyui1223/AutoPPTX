from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import sys
import math

prs = Presentation()
# adjustable, in inches
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6] # an empty slide

def get_image_size(image_path): # obtain size of image in inches
    img = Image.open(image_path)
    pixelWidth = float(img.size[0])
    pixelHeight = float(img.size[1])
    try:
        dpi = float(img.info['dpi'][0])
    except:
        #print("image "+image_path+" does not contain dpi info. uses dpi=96 for now")
        dpi = 96 # tekitou
    imageWidth = pixelWidth/dpi
    imageHeight = pixelHeight/dpi
    return imageWidth, imageHeight


def align_CTA_style(imageNames, mW, mH, slideWidth, slideHeight):
    columns={}
    # put each of sorted figureName lists to a dictionary
    # the name of order in this dictionaly decides the figure order in pptx.
    for TCA in ["CTL","TST","ANL","CMP","ETC","MAM","JJA","SON","DJF"]:      
        if any(iN[:3] == TCA for iN in imageNames): # only if figureName that start with specified header exists
            columns[TCA]=sorted([iN for iN in imageNames if TCA == iN[:3]])
        
    num_columns=0
    num_rows=0
    for TCA, figNameList in columns.items(): 
        if len(figNameList) > num_rows:
            num_rows=len(figNameList)
        num_columns+=1

    # transpose if it can make figure bigger
    asis=abs(slideHeight/slideWidth-num_rows*mH/(num_columns*mW))
    transpose=abs(slideHeight/slideWidth-num_columns*mH/(num_rows*mW))
    if asis > transpose:
        tp=True
    else:
        tp=False
    
    return get_positions(columns,num_columns,num_rows,mW, mH,slideWidth, slideHeight,tp)

        
def align_tile_style(imageNames, mW, mH, slideWidth, slideHeight):
    columns={}
    num_columns, num_rows = optimize(len(imageNames), mW, mH, slideWidth, slideHeight)
    i = 0
    imageNames=sorted(imageNames)
    for column_name in range(1, num_columns+1):
        columns[column_name]=imageNames[i:i+num_rows]
        i += num_rows
    return get_positions(columns,num_columns,num_rows,mW, mH,slideWidth, slideHeight, False)


def optimize(num_images, mW, mH, sW, sH):
    amari=99999
    if num_images == 1:
        return 1, 1
    for i in range(1,num_images+1):
        num_columns = i
        num_rows = math.ceil(num_images/i)
        if amari > abs(sH/sW-mH*num_rows/mW/num_columns):
          amari = abs(sH/sW-mH*num_rows/mW/num_columns)
          opt_columns = num_columns
          opt_rows = num_rows
    return opt_columns, opt_rows


# obtain the left, top, and height positions for each figures for the specified set of columns, rows
# to maximize the figure sizes while fitting into the slide-size-minus-padding.
def get_positions(columns, num_columns, num_rows, mW, mH, sW, sH, tp):
    if tp== True:
        figList, topList, leftList, ratio = get_positions(columns, num_rows, num_columns, mH, mW, sW, sH, False)
        return figList, leftList, topList, ratio
    else:  
        figList=[]
        topList=[]
        leftList=[]
        heightList=[]
        left=0
        ratio_w=sW/(num_columns*mW)
        ratio_h=sH/(num_rows*mH)
        ratio = min(ratio_w,ratio_h)

        for TCA, figNameList in columns.items(): 
            top=0
            for fig in figNameList:
                figList.append(fig)
                topList.append(top)
                leftList.append(left)
                top+=mH*ratio
            left+=mW*ratio
        return figList, topList, leftList, ratio


def align_images(image_path, mW, mH, slideWidth, slideHeight):
    imageNames=os.listdir(image_path)
    headers= ["CTL","TST","ANL","CMP","ETC","MAM","JJA","SON","DJF"]
    if all(iN[:3] in headers for iN in imageNames): # all imageNames starts from any of headers
        print("uses alignment for test, control comparison")
        figNameList, topList, leftList, ratio = align_CTA_style(imageNames, mW, mH, slideWidth, slideHeight)
    elif all(iN[:1] in [str(i) for i in range(10)] for iN in imageNames): # all imageNames starts from integer numbers
        print("uses tiling alignment")
        figNameList, topList, leftList, ratio = align_tile_style(imageNames, mW, mH, slideWidth, slideHeight)
    else:
        print("figure names are out of order!")
    return figNameList, topList, leftList, ratio


def obtain_max_image_size(image_path): # obtain maximum image size in n-th directory
    maxImageWidth=maxImageHeight=0
    for img in os.listdir(image_path):
        imageWidth, imageHeight = get_image_size(image_path+"/"+img)
        if maxImageWidth < imageWidth:   maxImageWidth=imageWidth 
        if maxImageHeight < imageHeight: maxImageHeight=imageHeight 
    return maxImageWidth, maxImageHeight


def get_image_height(image_path):
    imageWidth, imageHeight = get_image_size(image_path)
    return imageHeight


def obtain_padding_list(root_path, sd):
    default = 1 #inches
    try:
        with open(root_path+"/padding.txt") as f:
            for line in f:
                if sd in line:
                    paddingList = line.split()
                    # given strings like slideid( 1( 2 3 4))
                    paddingList.pop(0)
                    if len(paddingList) == 0:
                        return default,default,default,default
                    if len(paddingList) == 1:
                        padding=float(paddingList[0])
                        return padding,padding,padding,padding
                    if len(paddingList) == 4:
                        return [float(padding) for padding in paddingList]
                    else:
                        print("too many or too less number(s) for "+sd+" in padding.txt")
        # if padding.txt does not contain slide info
        return default,default,default,default
    except:
        print("Not found: "+root_path+"/padding.txt")
        return default,default,default,default


root_path = sys.argv[1]
for sd in sorted(os.listdir(root_path)):
    if "slide" not in sd:
        continue
        
    print("processing " +sd+ "...")
    slide = prs.slides.add_slide(blank_slide_layout)
    image_path = root_path+'/'+sd
    maxImageWidth, maxImageHeight = obtain_max_image_size(image_path)
    if max(maxImageWidth,maxImageHeight)==0:
      continue
    # obtain the space in the slide to be left empty. Better alert if bigger than slide itself
    paddingLeft, paddingRight, paddingTop, paddingBottom = obtain_padding_list(root_path, sd)
    figList, topList, leftList, ratio = align_images(image_path, maxImageWidth, maxImageHeight, 
                                                     prs.slide_width -Inches(paddingLeft + paddingRight), 
                                                     prs.slide_height-Inches(paddingTop + paddingBottom))
    for fig, top, left in zip(figList, topList, leftList):
        pic = slide.shapes.add_picture(image_path+"/"+fig,
                                       left+Inches(paddingLeft), top+Inches(paddingTop),
                                       height=ratio*get_image_height(image_path+"/"+fig))

prs.save(root_path+'/temp.pptx')


